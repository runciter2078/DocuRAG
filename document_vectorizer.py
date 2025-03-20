# -*- coding: utf-8 -*-
"""
############################# DocuRAG #############################
----------------------- Document Vectorizer Script -----------------------

Description:
This script monitors a specified folder for new files (PDF, DOCX, TXT, PPTX),
converts their content to text, generates embeddings using the 'LaBSE' SentenceTransformer model,
and stores these embeddings in a ChromaDB vector database. It also logs processed files
to avoid reprocessing.

Functionality:
1. Detects new files:
   - Scans the specified folder (data_path) for allowed file types.
   - Checks processed files using their names and hashes.
   - Only processes files that are new or have changed, excluding files with certain extensions.

2. Text Conversion:
   - Converts file content to plain text using libraries specific to each file type.

3. Embedding Generation:
   - Splits text into semantic blocks and generates embeddings using the 'LaBSE' model.

4. Storage in ChromaDB:
   - Stores embeddings in the ChromaDB database, avoiding duplicates.

5. Processed Files Logging:
   - Updates the 'processed_files.txt' file with the processed files' information.

Expected Output:
- New embeddings stored in ChromaDB.
- An updated 'processed_files.txt' log file.
- Console messages indicating progress and errors.
"""

import warnings
warnings.simplefilter(action='ignore', category=FutureWarning)

import os
import re
import fitz  # PyMuPDF for PDFs
import chromadb
import mimetypes
from sentence_transformers import SentenceTransformer
from docx import Document
from pptx import Presentation
from nltk.tokenize import sent_tokenize
import hashlib
import yaml
import nltk

def load_config(config_file='config.yml'):
    script_dir = os.path.dirname(os.path.abspath(__file__))
    config_path = os.path.join(script_dir, config_file)
    with open(config_path, 'r', encoding='utf-8') as f:
        config = yaml.safe_load(f)
    return config

config = load_config()

# Set BASE_PATH to a relative path for generic usage.
BASE_PATH = config.get('BASE_PATH', '.')
nltk_data_path = os.path.join(BASE_PATH, config.get('common', {}).get('nltk_data_path', 'nltk_data'))

nltk.data.path.append(nltk_data_path)
if not os.path.exists(nltk_data_path):
    os.makedirs(nltk_data_path)

try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    print(f"Downloading 'punkt' data in {nltk_data_path}...")
    nltk.download('punkt', download_dir=nltk_data_path)

dv_config = config.get('document_vectorizer', {})
DATA_PATH = os.path.join(BASE_PATH, dv_config.get('data_path', 'Data'))
DB_PATH = os.path.join(BASE_PATH, dv_config.get('db_path', 'sqlite_folder'))
PROCESSED_FILES_LOG = os.path.join(BASE_PATH, dv_config.get('processed_files_log', 'processed_files.txt'))
SENTENCE_TRANSFORMER_MODEL = dv_config.get('sentence_transformer_model', 'sentence-transformers/LaBSE')
EMBEDDING_BLOCK_SIZE = dv_config.get('embedding_block_size', 1000)
N_RESULTS = dv_config.get('n_results', 10)
EXCLUDE_EXTENSIONS = dv_config.get('exclude_extensions', [".xlsx", ".xls"])

def create_folder(path):
    if not os.path.exists(path):
        os.makedirs(path)

create_folder(DATA_PATH)
create_folder(DB_PATH)

sentence_model = SentenceTransformer(SENTENCE_TRANSFORMER_MODEL)

def load_processed_files(log_path):
    processed_files = {}
    if os.path.exists(log_path):
        try:
            with open(log_path, 'r', encoding='utf-8', errors='ignore') as f:
                for line in f:
                    if '\t' in line:
                        file_name, file_hash = line.strip().split('\t')
                        processed_files[file_name] = file_hash
        except Exception as e:
            print(f"Error reading the log file: {e}")
    return processed_files

def log_processed_file(log_path, file_name, file_hash):
    try:
        with open(log_path, 'a', encoding='utf-8') as f:
            f.write(f"{file_name}\t{file_hash}\n")
    except Exception as e:
        print(f"Error writing to the log file: {e}")

def calculate_file_hash(file_path):
    hash_md5 = hashlib.md5()
    try:
        with open(file_path, 'rb') as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
    except Exception as e:
        print(f"Error calculating hash for {file_path}: {e}")
    return hash_md5.hexdigest()

def get_mime_type(file_path):
    mime_type, _ = mimetypes.guess_type(file_path)
    return mime_type

def generate_embeddings(text_chunks, model, max_chunks=None):
    if max_chunks:
        text_chunks = text_chunks[:max_chunks]
    embeddings = model.encode(text_chunks)
    return embeddings

def convert_to_text(file_path):
    mime_type = get_mime_type(file_path)
    
    if mime_type == 'application/pdf':
        return extract_text_from_pdf(file_path)
    elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        return extract_text_from_docx(file_path)
    elif mime_type == 'text/plain':
        return extract_text_from_txt(file_path)
    elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        return extract_text_from_pptx(file_path)
    else:
        print(f"Unsupported file type: {mime_type}")
        return None

def extract_text_from_pdf(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        text = ""
        for page in doc:
            text += page.get_text("text")
        doc.close()
        return clean_text(text)
    except Exception as e:
        print(f"Error extracting text from PDF {pdf_path}: {e}")
        return ""

def extract_text_from_docx(docx_path):
    try:
        doc = Document(docx_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return clean_text(text)
    except Exception as e:
        print(f"Error extracting text from DOCX {docx_path}: {e}")
        return ""

def extract_text_from_txt(txt_path):
    try:
        with open(txt_path, 'r', encoding='utf-8') as f:
            text = f.read()
        return clean_text(text)
    except UnicodeDecodeError:
        try:
            with open(txt_path, 'r', encoding='latin-1') as f:
                text = f.read()
            return clean_text(text)
        except Exception as e:
            print(f"Error extracting text from TXT {txt_path}: {e}")
            return ""
    except Exception as e:
        print(f"Error extracting text from TXT {txt_path}: {e}")
        return ""

def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + ' '
                if hasattr(slide, 'notes_slide') and slide.notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    text += notes + ' '
        return clean_text(text)
    except Exception as e:
        print(f"Error extracting text from PPTX {pptx_path}: {e}")
        return ""

def clean_text(text):
    text = re.sub(r'\n+', ' ', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text

def split_text_into_chunks(text, max_length=EMBEDDING_BLOCK_SIZE):
    sentences = sent_tokenize(text, language='spanish')
    chunks = []
    current_chunk = ''
    for sentence in sentences:
        if len(current_chunk) + len(sentence) <= max_length:
            current_chunk += ' ' + sentence
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = sentence
    if current_chunk:
        chunks.append(current_chunk.strip())
    return chunks

def store_embeddings_in_chroma(embeddings, text_chunks, doc_name):
    try:
        client = chromadb.PersistentClient(path=DB_PATH)
        collection = client.get_or_create_collection(name="document_embeddings")
    except Exception as e:
        print(f"Error connecting to ChromaDB: {e}")
        return

    embeddings = [embedding.tolist() for embedding in embeddings]
    ids = [f"{doc_name}_chunk_{i}" for i in range(len(text_chunks))]
    metadatas = [{"source": doc_name} for _ in range(len(text_chunks))]

    existing_ids = set()
    try:
        existing_docs = collection.get(ids=ids)
        existing_ids = set(existing_docs["ids"])
    except Exception as e:
        print(f"Error retrieving existing documents from ChromaDB: {e}")

    new_embeddings, new_documents, new_metadatas, new_ids = [], [], [], []
    for i, embedding_id in enumerate(ids):
        if embedding_id not in existing_ids:
            new_embeddings.append(embeddings[i])
            new_documents.append(text_chunks[i])
            new_metadatas.append(metadatas[i])
            new_ids.append(embedding_id)

    if new_embeddings:
        try:
            collection.add(
                embeddings=new_embeddings,
                documents=new_documents,
                metadatas=new_metadatas,
                ids=new_ids
            )
            print(f"{len(new_embeddings)} new embeddings stored in ChromaDB for {doc_name}.")
        except Exception as e:
            print(f"Error adding embeddings to ChromaDB: {e}")

def main():
    processed_files = load_processed_files(PROCESSED_FILES_LOG)
    try:
        files = [f for f in os.listdir(DATA_PATH) if os.path.isfile(os.path.join(DATA_PATH, f))]
    except FileNotFoundError:
        print(f"DATA_PATH folder does not exist: {DATA_PATH}")
        return

    for file_name in files:
        if not file_name.startswith('~$'):
            file_extension = os.path.splitext(file_name)[1].lower()
            if file_extension in EXCLUDE_EXTENSIONS:
                print(f"Excluding file due to extension: {file_name}")
                continue

            file_path = os.path.join(DATA_PATH, file_name)
            file_hash = calculate_file_hash(file_path)

            if file_name not in processed_files or processed_files[file_name] != file_hash:
                print(f"Processing file: {file_name}")
                text = convert_to_text(file_path)
                if text:
                    text_chunks = split_text_into_chunks(text)
                    embeddings = generate_embeddings(text_chunks, sentence_model)
                    store_embeddings_in_chroma(embeddings, text_chunks, file_name)

                processed_files[file_name] = file_hash
                log_processed_file(PROCESSED_FILES_LOG, file_name, file_hash)
            else:
                print(f"File {file_name} already processed and unchanged.")

    print("Processing complete.")

if __name__ == "__main__":
    main()
